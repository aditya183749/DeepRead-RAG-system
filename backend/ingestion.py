"""
ingestion.py — Full document processing pipeline.

Pipeline:
  1. Detect file type (PDF / DOCX)
  2. Extract text page-by-page with section heading detection
  3. Split text into overlapping chunks (300 tokens, 30 overlap)
  4. Extract entities per chunk for rich metadata
  5. Generate BGE-M3 embeddings via sentence-transformers (local, no Ollama)
  6. Add to FAISS + BM25 indexes

All metadata is persisted alongside each chunk so it NEVER vanishes.
"""

import uuid
import re
import datetime
import fitz                                     # PyMuPDF — PDF parsing
from docx import Document as DocxDocument       # python-docx — DOCX parsing
from pathlib import Path
from typing import List, Tuple, Dict, Any, Callable, Optional
import csv
import io

import numpy as np
import tiktoken
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer

from config import (
    CHUNK_SIZE, CHUNK_OVERLAP, EMBED_MODEL, EMBED_DIM
)
import vector_store
import bm25_store


# ─── Embedding Model (singleton — load once, reuse) ───────────────────────────

_embed_model: SentenceTransformer | None = None


def get_embed_model() -> SentenceTransformer:
    """
    Lazy-load BGE-M3 on first call.
    First call downloads ~2.2 GB weights (cached locally after that).
    normalize_embeddings=True ensures L2-unit vectors → cosine = inner product.
    """
    global _embed_model
    if _embed_model is None:
        print(f"[Embeddings] Loading {EMBED_MODEL} (may download on first run)...")
        _embed_model = SentenceTransformer(EMBED_MODEL)
        print(f"[Embeddings] {EMBED_MODEL} ready. Dim={EMBED_DIM}")
    return _embed_model


# ─── Tokenizer (for chunk sizing) ─────────────────────────────────────────────

_tokenizer = tiktoken.get_encoding("cl100k_base")


def _token_length(text: str) -> int:
    return len(_tokenizer.encode(text))


# ─── Text Splitter ────────────────────────────────────────────────────────────

def _get_splitter() -> RecursiveCharacterTextSplitter:
    """
    RecursiveCharacterTextSplitter with token-accurate sizing.
    Separator priority: paragraph → line → sentence → word.
    chunk_size=300 tokens, overlap=30 tokens (team decision).
    """
    return RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " "],
        length_function=_token_length,
        is_separator_regex=False
    )


# ─── Entity Extraction ────────────────────────────────────────────────────────

def _extract_entities(text: str) -> Dict[str, List[str]]:
    """
    Lightweight regex-based entity extraction for chunk metadata enrichment.

    Extracts:
      - metrics    : numbers with units/% e.g. "98.2% accuracy", "$4.2B"
      - variables  : defined terms e.g. "α = 0.01", "N = 512"
      - dates      : years and date patterns e.g. "2023", "March 2021"
      - emails     : email addresses

    NOTE: Full NER (spaCy) will be added in Phase 3.
          This lightweight version runs with zero extra dependencies.
    """
    metrics = re.findall(
        r'\b\d+\.?\d*\s*(%|billion|million|trillion|accuracy|F1|BLEU|'
        r'tokens|ms|GB|TB|KB|MB|fps|GHz|MHz|kB|ROUGE)\b',
        text, re.IGNORECASE
    )
    variables = re.findall(
        r'\b([A-Za-z_]\w*)\s*[=:]\s*[\d\.]+\b',
        text
    )
    dates = re.findall(
        r'\b((?:Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|'
        r'Jun(?:e)?|Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|'
        r'Nov(?:ember)?|Dec(?:ember)?)\s+\d{4}|\d{4})\b',
        text
    )
    emails = re.findall(r'\b[\w.+-]+@[\w-]+\.[a-zA-Z]{2,}\b', text)

    return {
        "metrics":   [m[0] if isinstance(m, tuple) else m for m in metrics][:10],
        "variables": list(set(variables))[:10],
        "dates":     list(set(dates))[:5],
        "emails":    emails[:5]
    }


# ─── PDF Extraction ───────────────────────────────────────────────────────────

def extract_text_from_pdf(file_path: Path) -> List[Dict[str, Any]]:
    """
    Extract text from a PDF, page by page, with section heading detection.

    Returns list of page dicts:
      [{"page_number": 1, "text": "...", "heading": "Introduction"}, ...]

    Heading detection: lines that are ALL-CAPS or short (<= 60 chars) and
    followed by a blank line are treated as section headings.
    """
    pages = []
    doc = fitz.open(str(file_path))

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text").strip()
        if not text:
            continue

        # Detect the first plausible heading on this page
        heading = _detect_heading(text)

        pages.append({
            "page_number": page_num + 1,
            "text":        text,
            "heading":     heading
        })

    doc.close()
    return pages


def _detect_heading(text: str) -> str:
    """Heuristic: first non-empty, short line is likely a section heading."""
    for line in text.split("\n"):
        line = line.strip()
        if line and len(line) <= 80 and not line.endswith(","):
            return line
    return ""


# ─── DOCX Extraction ──────────────────────────────────────────────────────────

def extract_text_from_docx(file_path: Path) -> List[Dict[str, Any]]:
    """
    Extract text from a DOCX, grouping every 3000 chars as a 'page'.
    Captures paragraph style names to identify headings.
    """
    doc = DocxDocument(str(file_path))

    full_text = ""
    current_heading = ""

    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        if para.style.name.startswith("Heading"):
            current_heading = t
        full_text += t + "\n"

    return _text_to_pages(full_text, current_heading)


# ─── TXT / Markdown Extraction ────────────────────────────────────────────────

def extract_text_from_txt(file_path: Path) -> List[Dict[str, Any]]:
    """Plain text and Markdown — read as UTF-8, split into 3000-char pages."""
    text = file_path.read_text(encoding="utf-8", errors="replace")
    heading = _detect_heading(text)
    return _text_to_pages(text, heading)


# ─── CSV Extraction ───────────────────────────────────────────────────────────

def extract_text_from_csv(file_path: Path) -> List[Dict[str, Any]]:
    """
    Convert CSV rows to readable text blocks.
    Each row becomes: "col1: val1 | col2: val2 | …"
    Groups ~50 rows per page.
    """
    rows = []
    with open(file_path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.DictReader(f)
        headers = reader.fieldnames or []
        for row in reader:
            line = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
            if line:
                rows.append(line)

    ROWS_PER_PAGE = 50
    pages = []
    for i, start in enumerate(range(0, len(rows), ROWS_PER_PAGE)):
        block = "\n".join(rows[start:start + ROWS_PER_PAGE])
        pages.append({"page_number": i + 1, "text": block, "heading": ", ".join(headers[:5])})
    return pages


# ─── XLSX Extraction ──────────────────────────────────────────────────────────

def extract_text_from_xlsx(file_path: Path) -> List[Dict[str, Any]]:
    """
    Read all sheets from an Excel file.
    Each sheet becomes one or more pages (50 rows per page).
    """
    import openpyxl
    wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
    pages = []
    page_num = 1
    for sheet in wb.worksheets:
        rows = []
        headers = []
        for r_idx, row in enumerate(sheet.iter_rows(values_only=True)):
            cells = [str(c) if c is not None else "" for c in row]
            if r_idx == 0:
                headers = [c for c in cells if c]
                continue
            line = " | ".join(f"{h}: {v}" for h, v in zip(headers, cells) if v)
            if line:
                rows.append(line)

        ROWS_PER_PAGE = 50
        for start in range(0, max(len(rows), 1), ROWS_PER_PAGE):
            block = "\n".join(rows[start:start + ROWS_PER_PAGE])
            if block:
                pages.append({"page_number": page_num, "text": block,
                               "heading": f"Sheet: {sheet.title}"})
                page_num += 1
    wb.close()
    return pages


# ─── PPTX Extraction ──────────────────────────────────────────────────────────

def extract_text_from_pptx(file_path: Path) -> List[Dict[str, Any]]:
    """
    Extract text from PowerPoint. Each slide = one page.
    Slide title becomes the heading.
    """
    from pptx import Presentation
    prs = Presentation(str(file_path))
    pages = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        heading = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if not t:
                    continue
                if not heading:   # first text block = slide title
                    heading = t
                texts.append(t)
        block = "\n".join(texts)
        if block:
            pages.append({"page_number": slide_num, "text": block, "heading": heading})
    return pages


# ─── HTML Extraction ──────────────────────────────────────────────────────────

def extract_text_from_html(file_path: Path) -> List[Dict[str, Any]]:
    """Strip HTML tags and extract visible text into 3000-char pages."""
    from bs4 import BeautifulSoup
    html = file_path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")
    # Remove scripts and styles
    for tag in soup(["script", "style", "head", "meta"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    # Collapse blank lines
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    heading = _detect_heading(text)
    return _text_to_pages(text, heading)


# ─── Shared Helper ────────────────────────────────────────────────────────────

def _text_to_pages(text: str, heading: str = "", page_size: int = 3000) -> List[Dict[str, Any]]:
    """Split a flat text string into page-sized dicts for the chunker."""
    pages = []
    for i, start in enumerate(range(0, len(text), page_size)):
        chunk = text[start:start + page_size].strip()
        if chunk:
            pages.append({
                "page_number": i + 1,
                "text":        chunk,
                "heading":     heading if i == 0 else ""
            })
    return pages


# ─── Chunk Pages ──────────────────────────────────────────────────────────────

def chunk_pages(
    pages: List[Dict[str, Any]],
    source_id: str,
    filename: str
) -> List[Dict[str, Any]]:
    """
    Split each page's text into overlapping chunks with full metadata.

    Every chunk carries complete provenance so it NEVER loses context:
      - source_id, filename, version
      - page_number, chunk_index, total_pages
      - section_heading (nearest heading)
      - chunk_type, chunk_strategy, embed_model
      - entity extraction results
      - ingested_at timestamp
      - token_count
    """
    splitter = _get_splitter()
    chunks = []
    chunk_index = 0
    total_pages = len(pages)
    ingested_at = datetime.datetime.utcnow().isoformat() + "Z"

    for page in pages:
        page_chunks = splitter.split_text(page["text"])

        for chunk_text in page_chunks:
            entities = _extract_entities(chunk_text)

            chunks.append({
                "id":   f"{source_id}_{chunk_index}",
                "text": chunk_text,
                "metadata": {
                    # ── Identity ────────────────────────────────────────────
                    "source_id":     source_id,
                    "filename":      filename,
                    "version":       1,

                    # ── Location ────────────────────────────────────────────
                    "page_number":   page["page_number"],
                    "chunk_index":   chunk_index,
                    "total_pages":   total_pages,
                    "section_heading": page.get("heading", ""),

                    # ── Type ───────────────────────────────────────────────
                    "chunk_type":     "text",
                    "chunk_strategy": "recursive",

                    # ── Embeddings ─────────────────────────────────────────
                    "embed_model":   EMBED_MODEL,
                    "embed_dim":     EMBED_DIM,

                    # ── Entities (metrics & variables from this chunk) ─────
                    "metrics":       entities["metrics"],
                    "variables":     entities["variables"],
                    "dates":         entities["dates"],

                    # ── Time ───────────────────────────────────────────────
                    "ingested_at":   ingested_at,
                    "content_date":  "",       # filled by user / future NER

                    # ── Quality ────────────────────────────────────────────
                    "token_count":   _token_length(chunk_text),
                    "retrieval_count": 0,
                }
            })
            chunk_index += 1

    return chunks


# ─── Generate Embeddings (BGE-M3) ─────────────────────────────────────────────

def generate_embeddings(texts: List[str]) -> List[List[float]]:
    """
    Generate L2-normalised 1024-dim BGE-M3 embeddings via sentence-transformers.

    Runs locally — no Ollama, no API call.
    batch_size=32 balances speed vs RAM on typical laptops.

    CRITICAL: same model MUST be used for both ingest and query embedding.
    'embed_model' stored in chunk metadata enforces this.
    """
    model = get_embed_model()
    embeddings = model.encode(
        texts,
        normalize_embeddings=True,  # L2-unit vectors → cosine = inner product
        batch_size=32,
        show_progress_bar=len(texts) > 50
    )
    return embeddings.tolist()


# ─── Main Ingestion Pipeline ───────────────────────────────────────────────────

def ingest_document(
    file_path: Path,
    filename: str,
    progress_cb=None          # optional callable(step: str, detail: str)
) -> Tuple[str, int, int]:
    """
    Full pipeline: file → chunks → embeddings → FAISS + BM25.

    Steps:
      1. Generate UUID as source_id
      2. Detect file type, extract text page-by-page
      3. Split into overlapping chunks with full metadata
      4. Generate BGE-M3 embeddings (batch)
      5. Add to FAISS index (vector search)
      6. Add to BM25 corpus (keyword search)

    Returns:
      (source_id, chunk_count, page_count)
    """
    def _progress(step: str, detail: str = ""):
        if progress_cb:
            progress_cb(step, detail)

    source_id = str(uuid.uuid4())

    # Step 2: Extract text
    _progress("extracting", f"Reading '{filename}'…")
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        pages = extract_text_from_pdf(file_path)
    elif suffix in (".docx", ".doc"):
        pages = extract_text_from_docx(file_path)
    elif suffix in (".txt", ".md"):
        pages = extract_text_from_txt(file_path)
    elif suffix == ".csv":
        pages = extract_text_from_csv(file_path)
    elif suffix in (".xlsx", ".xls"):
        pages = extract_text_from_xlsx(file_path)
    elif suffix == ".pptx":
        pages = extract_text_from_pptx(file_path)
    elif suffix in (".html", ".htm"):
        pages = extract_text_from_html(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")

    if not pages:
        raise ValueError("No text could be extracted from the document.")

    _progress("extracted", f"Extracted {len(pages)} pages")

    # Step 3: Chunk with rich metadata
    _progress("chunking", f"Splitting {len(pages)} pages into chunks…")
    chunks = chunk_pages(pages, source_id, filename)
    if not chunks:
        raise ValueError("Document produced no chunks after splitting.")

    _progress("chunked", f"Created {len(chunks)} chunks (300 tokens each, 30 overlap)")
    print(f"[Ingest] '{filename}' → {len(pages)} pages, {len(chunks)} chunks")

    # Step 4: Embed all chunks in batch
    _progress("embedding", f"Generating BGE-M3 embeddings for {len(chunks)} chunks…")
    texts = [c["text"] for c in chunks]
    embeddings = generate_embeddings(texts)
    _progress("embedded", f"Embeddings done ({len(embeddings[0])}-dim vectors)")

    # Step 5: Add to FAISS
    _progress("indexing_faiss", "Adding to FAISS vector index…")
    vector_store.add_chunks(
        ids=[c["id"] for c in chunks],
        embeddings=embeddings,
        documents=texts,
        metadatas=[c["metadata"] for c in chunks]
    )
    _progress("indexed_faiss", "FAISS index updated")

    # Step 6: Add to BM25
    _progress("indexing_bm25", "Adding to BM25 keyword index…")
    bm25_store.add_chunks(
        ids=[c["id"] for c in chunks],
        documents=texts,
        metadatas=[c["metadata"] for c in chunks]
    )
    _progress("done", f"'{filename}' fully indexed — {len(chunks)} chunks ready")

    print(f"[Ingest] '{filename}' indexed. source_id={source_id}")
    return source_id, len(chunks), len(pages)
